#%% 
from ani_tools import *
import os
import re
from datetime import datetime
import subprocess
from pptx import Presentation
from ppt2video.tools import _clean_text

CONF_FILE = '../config/config.json'

def find_pptx_files(base_dir):
    category_K_files = []
    category_E_files = []
    
    # Walk through the directory and its subdirectories
    for root, _, files in os.walk(base_dir):
        for file in files:
            if file.endswith('.pptx'):
                if '_K_' in file:
                    category_K_files.append(os.path.join(root, file))
                elif '_E_' in file:
                    category_E_files.append(os.path.join(root, file))
    
    return category_K_files, category_E_files

def filter_long_files(files): # non 13sec shorts
    return [file for file in files if '_shorts' not in file.lower()]

def filter_short_files(files): # non 13sec shorts
    return [file for file in files if '_shorts' in file.lower() and '_13sec' not in file.lower()]

def filter_13sec_short_files(files):
    return [file for file in files if '_13sec' in file.lower()]

def sort_files_by_date(file_list):
    files_with_date = []
    files_no_date = []
    
    for file in file_list:
        match = re.search(r'\d{4}-\d{2}-\d{2}', file)
        if match:
            date_str = match.group(0)
            # Convert the found date string into a datetime object
            date_obj = datetime.strptime(date_str, '%Y-%m-%d')
            files_with_date.append((file, date_obj))
        else: 
            files_no_date.append(file)
    
    if len(files_no_date) > 0: 
        raise Exception('Certain files do not have proper date in filename')

    # Sort the files by date (oldest to newest)
    files_with_date.sort(key=lambda x: x[1])
    return [file[0] for file in files_with_date]

def lprint(flist):
    for f in flist:
        print(f)

def open_ppt_file(ppt_path):
    try:
        subprocess.run(['start', 'powerpnt', ppt_path], check=True, shell=True)
    except subprocess.CalledProcessError as e:
        print(f"Error occurred while opening the file: {e}")

def get_slide_count(ppt_path):
    ppt = Presentation(ppt_path)
    return len(ppt.slides)

def get_notes(ppt_path):
    ppt = Presentation(ppt_path)

    slide_content = ''
    for slide_number, slide in enumerate(ppt.slides):
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            notes = _clean_text(notes)
            slide_content += f'<br> (p{slide_number+1}) '+ notes + '\n'
    
    return slide_content

def set_notes(ppt_path, text):
    notes = text.splitlines()
    pattern = r"<br> \(p\d+\) "
    ppt = Presentation(ppt_path)

    for slide_number, slide in enumerate(ppt.slides):
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            slide.notes_slide.notes_text_frame.text = re.sub(pattern, "", notes[slide_number])

    ppt.save(ppt_path)

def get_prefix_from_filename(file):
    filename = os.path.basename(file)
    if '13sec' in filename: 
        prefix = '[13 sec'
    elif 'shorts' in filename:
        prefix = '[Shorts'
    else:    
        prefix = '[Video'
    date_format = r"\b\d{4}-\d{2}-\d{2}\b"
    match = re.search(date_format, filename)
    date = match.group(0) if match else ''
    if date != '':
        prefix = prefix + ' ' + date + '] '
    else: 
        prefix = prefix + '] '
    return prefix

#%% 
base_dir = r'C:\Users\user\projects\analysis'
cat_K_files, cat_E_files = find_pptx_files(base_dir)
cat_K_longs = sort_files_by_date(filter_long_files(cat_K_files))
cat_E_longs = sort_files_by_date(filter_long_files(cat_E_files))
cat_K_shorts = sort_files_by_date(filter_short_files(cat_K_files))
cat_E_shorts = sort_files_by_date(filter_short_files(cat_E_files))
cat_K_13secs = sort_files_by_date(filter_13sec_short_files(cat_K_files))
cat_E_13secs = sort_files_by_date(filter_13sec_short_files(cat_E_files))

display(cat_K_longs)
display(cat_K_shorts)
display(cat_K_13secs)
display(cat_E_longs)
display(cat_E_shorts)
display(cat_E_13secs)

#%% 
K_file = cat_K_shorts[]
E_file = K_file.replace('_K_', '_E_')
K_note = get_notes(K_file)
E_note = translate_script(K_note, CONF_FILE)
set_notes(E_file, E_note) 
get_desc(K_note, CONF_FILE)

#%% 

QP_SHORTS_CARD = 45  # Quarterly Performances -Shorts (한글본, Korean)
ist = IST()

form_data = {
    "title": "Programmatically Created Post",
    "content": "",
    "tags": "haha this is tags", 
}
images = []

ist.create_post(QP_SHORTS_CARD, form_data, images, html=True)



