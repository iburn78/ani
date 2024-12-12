#%% 
import requests
import json
from bs4 import BeautifulSoup
from urllib.parse import urljoin
import os
import re
from datetime import datetime, date
import subprocess
from pptx import Presentation
import win32com.client
from ppt2video.tools import _clean_text
from ani_tools import *
import shutil

CONF_FILE = '../config/config.json'
BASE_DIR = r'C:\Users\user\projects\analysis'
PPT_WORK_DIR = r'C:\Users\user\projects\ani\data\ppt'
# Better to clean the slide image root directory before running...
SLIDE_IMAGE_ROOT = 'C:\\Users\\user\\projects\\analysis\\temp\\'
MAX_IMAGES_PER_POST = 10

QP_SHORTS_Korean = 45  # Quarterly Performances -Shorts (한글본, Korean)
QP_SHORTS_English = 50  # Quarterly Performances -Shorts (English)
QP_13sec_Korean = 51  # Quarterly Performances -13sec (한글본, Korean)
QP_13sec_English = 52  # Quarterly Performances -13sec (English)
QP_LF_Korean = 53  # Quarterly Performances -Long form (한글본, Korean)
QP_LF_English = 54  # Quarterly Performances -Long form (English)

Card_dict = {'E':{
    'LF': QP_LF_English, 
    'ST': QP_SHORTS_English,
    '13': QP_13sec_English,
}, 'K':{
    'LF': QP_LF_Korean, 
    'ST': QP_SHORTS_Korean,
    '13': QP_13sec_Korean,
}}

class IST:  # IssueTracker Handler
    IST_SITE = 'https://issuetracker.info/'

    def __init__(self, conf_file):
        self.session = requests.Session()
        self.CONF_FILE = conf_file
        self.login()

    def _get_csrf_token(self, url):
        response = self.session.get(url)
        soup = BeautifulSoup(response.text, "html.parser")
        csrf_token_input = soup.find("input", {"name": "csrfmiddlewaretoken"})
        self.csrf_token = csrf_token_input["value"] if csrf_token_input else None

    def login(self):
        login_url = urljoin(IST.IST_SITE, 'login/')
        self._get_csrf_token(login_url)

        with open(self.CONF_FILE, 'r') as json_file:
            config = json.load(json_file)
            self.ist_id = config['issuetracker_id']
            ist_pass = config['issuetracker_pass']

        form_data = {
            "username": self.ist_id,
            "password": ist_pass,
            "csrfmiddlewaretoken": self.csrf_token,
        }
        headers = {
            "Referer": login_url,  # Required by CSRF protection
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/117.0.0.0 Safari/537.36",
        }

        response = self.session.post(login_url, data=form_data, headers=headers)
        if response.status_code != 200:
            raise Exception('IST login failed')


    def create_post(self, card_id: int, form_data: dict, images: list = [], html=False):
        target_url = urljoin(IST.IST_SITE, f'card/{card_id}/new_post/')
        self._get_csrf_token(target_url)

        form_data['csrfmiddlewaretoken'] = self.csrf_token
        lm = len(images)
        if lm > 10: 
            raise Exception('Only up to 10 images can be posted')
        mimage_str = 'abcdefghij'
        form_data['mimage_keys'] = mimage_str[:lm].upper() + mimage_str[lm:]
        if html:
            form_data['html_or_text'] = 'html'

        files = {}
        for i, im in enumerate(images):
            if not os.path.exists(im):
                raise Exception(f"image {im} does not exist")
            files[f'image{i+1}_input'] = (os.path.basename(im), open(im, 'rb'))

        headers = {
            "Referer": target_url,  # Required by CSRF protection
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/117.0.0.0 Safari/537.36",
        }
        response = self.session.post(target_url, data=form_data, files=files, headers=headers)

        # Close the image files after the POST request
        for _,file in files.items():
            file[1].close()

        if response.status_code == 200:
            print(f"Post{' ['+form_data['title']+']' if 'title' in form_data else ''}{' ('+str(lm)+' images)' if lm>0 else ''} created.")
            return None
        else:
            print(f"Post{' ['+form_data['title']+']' if 'title' in form_data else ''}{' ('+str(lm)+' images)' if lm>0 else ''} creation FAILED.")
            return response.text



def find_pptx_files(base_dir):
    category_K_files = []
    category_E_files = []
    
    # Walk through the directory and its subdirectories
    for root, _, files in os.walk(base_dir):
        for file in sorted(files):
            if file.endswith('.pptx'):
                if '_K_' in file:
                    category_K_files.append(os.path.join(root, file))
                elif '_E_' in file:
                    category_E_files.append(os.path.join(root, file))
    
    return category_K_files, category_E_files

def filter_long_files(files): # non 13sec shorts
    return [file for file in files if '_shorts' not in file.lower()]

def filter_short_files(files): # non 13sec shorts
    return [file for file in files if '_shorts' in file.lower() and '_13sec' not in file.lower()]

def filter_13sec_short_files(files):
    return [file for file in files if '_13sec' in file.lower()]


def sort_files_by_date(file_list, dates_on_after=None):
    # Convert string input to date if provided, default to today's date
    if isinstance(dates_on_after, str):
        dates_on_after = datetime.strptime(dates_on_after, '%Y-%m-%d').date()
    elif dates_on_after is None:
        dates_on_after = date.today()

    files_with_date = []
    files_no_date = []

    for file in file_list:
        match = re.search(r'\d{4}-\d{2}-\d{2}', file)
        if match:
            date_str = match.group(0)
            # Convert the found date string into a datetime.date object
            date_obj = datetime.strptime(date_str, '%Y-%m-%d').date()
            if date_obj >= dates_on_after:
                files_with_date.append((file, date_obj))
        else: 
            files_no_date.append(file)
    
    if len(files_no_date) > 0: 
        raise Exception('Certain files do not have a proper date in the filename')

    # Sort the files by date (oldest to newest)
    files_with_date.sort(key=lambda x: x[1])
    return [file[0] for file in files_with_date]

# Easily open powerpoint software from python
def open_ppt_file(ppt_path):
    try:
        subprocess.run(['start', 'powerpnt', ppt_path], check=True, shell=True)
    except subprocess.CalledProcessError as e:
        print(f"Error occurred while opening the file: {e}")

def get_slide_count(ppt_path):
    ppt = Presentation(ppt_path)
    return len(ppt.slides)

def get_notes(ppt_path):
    ppt = Presentation(ppt_path)
    slide_content = ''
    for slide_number, slide in enumerate(ppt.slides):
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            notes = _clean_text(notes)
            slide_content += f'<br>(p{slide_number+1}) '+ notes + '\n'
    
    return slide_content

def set_notes(ppt_path, text):
    notes = text.splitlines()
    pattern = r"<br>\(p\d+\) "
    ppt = Presentation(ppt_path)

    if len(notes) != len(ppt.slides):
        raise Exception('Set Note Error')

    for slide_number, slide in enumerate(ppt.slides):
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            slide.notes_slide.notes_text_frame.text = re.sub(pattern, "", notes[slide_number])

    ppt.save(ppt_path)

def ppt_to_images(ppt_path, slide_image_root=SLIDE_IMAGE_ROOT):
    slide_folder = os.path.abspath(os.path.join(slide_image_root, os.path.basename(ppt_path).replace('.pptx', '')))
    os.makedirs(slide_folder, exist_ok=True)

    # Initialize PowerPoint
    ppt_app = win32com.client.Dispatch("PowerPoint.Application")
    presentation = ppt_app.Presentations.Open(ppt_path, WithWindow=False)

    # Loop through slides and save each as an image
    for i, slide in enumerate(presentation.Slides):
        image_path = os.path.join(slide_folder, f'slide{i}.png')
        slide.Export(image_path, "PNG")
        # print(f"Saved slide {i} to {image_path}")

    total_slides = len(presentation.Slides)

    # Close the presentation
    presentation.Close()
    ppt_app.Quit()
    return slide_folder, total_slides

def _type_checker(ppt_path, type):
    if '_K_' in os.path.basename(ppt_path): 
        lang = 'K'
    elif '_E_' in os.path.basename(ppt_path): 
        lang = 'E'
    else: 
        raise Exception('Check file name')
    
    if '_shorts' not in ppt_path and type == 'ST':
        raise Exception('Check type')

    if '_13sec' not in ppt_path and type == '13':
        raise Exception('Check type')

    if '_shorts' in ppt_path or '_13sec' in ppt_path:
        if type == 'LF':
            raise Exception('Check type')

    if type not in ['LF', 'ST', '13']:
        raise Exception('Check type')
    
    return lang

# type = 'LF', 'ST', '13'
def upload_a_ppt_ist(ist, ppt_path, type):
    youtube_log = pd.read_excel(YOUTUBE_LOG)
    vid_obj = youtube_log.loc[youtube_log['filename']==os.path.basename(ppt_path), 'id']

    if len(vid_obj) == 1: 
        vlink = f'https://youtu.be/{vid_obj.values[0]}'
        vtitle = youtube_log.loc[youtube_log['filename']==os.path.basename(ppt_path), 'title'].values[0].strip()
        vdesc = youtube_log.loc[youtube_log['filename']==os.path.basename(ppt_path), 'desc'].values[0].strip()
        vtags = youtube_log.loc[youtube_log['filename']==os.path.basename(ppt_path), 'keywords'].values[0].strip()
    else: 
        vlink = ''
        raise Exception('Check Video Link')

    lang = _type_checker(ppt_path, type) 
    if lang == '_E_': 
        vl_name = 'Issue Tracker'
    else: 
        vl_name = 'Quarterly Performances'
    card_id = Card_dict[lang][type]
    notes = get_notes(ppt_path)
    notes_split = notes.splitlines() # no new line included "\n"
    slide_folder, total_slides = ppt_to_images(ppt_path)
    images = sorted(
    [os.path.join(slide_folder, f) for f in os.listdir(slide_folder) if os.path.isfile(os.path.join(slide_folder, f))],
    key=lambda f: os.path.getctime(os.path.join(slide_folder, f))
    )

    total_files = (total_slides // MAX_IMAGES_PER_POST) + (1 if total_slides % MAX_IMAGES_PER_POST > 0 else 0)
    # [title, tags, desc] = get_desc(notes, lang, conf_file)
    title = vtitle
    tags = vtags.replace(',', '')

    desc_lines = vdesc.splitlines()
    if desc_lines and re.match(r'^(#\w+(\s+#\w+)*)$', desc_lines[0].strip()): 
        desc = "\n".join(desc_lines[1:])
    else: 
        raise Exception('Check desc from Excel')

    file_number = 1
    for start_idx in range(0, total_slides, MAX_IMAGES_PER_POST):
        notes_for_subfile = ''
        images_for_subfile = []
        for slide_idx in range(start_idx, min(start_idx + MAX_IMAGES_PER_POST, total_slides)):
            note_for_slide = notes_split[slide_idx]
            match = re.match(r"<br>\(p\d+\)\s*(.*)", note_for_slide)
            if match: 
                note_content = match.group(1).strip()
                if note_content: 
                    notes_for_subfile += note_for_slide + '\n'
            images_for_subfile.append(images[slide_idx])
        if total_files > 1:
            title_for_subfile = f'[{file_number}/{total_files}] '+ title
            desc_for_subfile = f'[{file_number}/{total_files}] '+ desc
            file_number += 1
        else: 
            title_for_subfile = title
            desc_for_subfile = desc
        if notes_for_subfile.startswith('<br>'):
            notes_for_subfile = notes_for_subfile[4:]
        content = desc_for_subfile + "\n<hr>" + notes_for_subfile
        content = content + f"<br><a href='{vlink}' target='_blank' rel='noopener noreferrer'>Youtube Link - {vl_name}</a>"

        form_data = {
            "title": title_for_subfile, 
            "content": content,
            "tags": tags.replace('#',''),
        }
        ist.create_post(card_id, form_data, images_for_subfile, html=True)


# If already not exists, creates a English copy of a Korean ppt file, and set notes with translated scripts
def set_translated_notes(K_file, conf_file):
    E_file = K_file.replace('_K_', '_E_')
    if not os.path.exists(E_file):
        shutil.copy(K_file, E_file)
    notes = get_notes(K_file)
    E_notes = translate_script(notes, conf_file)
    set_notes(E_file, E_notes)
    print(f'English copy of {os.path.basename(K_file)} completed with English notes')
    return None

def trans_list_of_K_files(K_list, E_list, conf_file):
    E_list_check = E_list.copy()
    for kf in K_list: 
        if kf.replace('_K_', '_E_') in E_list:
            set_translated_notes(kf, conf_file)
            E_list_check.remove(kf.replace('_K_', '_E_'))
        else: 
            raise Exception('Check creation of Engfile')
    if len(E_list_check) > 0: 
        raise Exception('Check creation of Engfile')


# Script translation might be different from youtube video, as it is done separately here....
def find_ppt_tranlate_and_upload(target_dir, dates_on_after, conf_file):
    cat_K_files, cat_E_files = find_pptx_files(target_dir)
    cat_K_longs = sort_files_by_date(filter_long_files(cat_K_files), dates_on_after)
    cat_E_longs = sort_files_by_date(filter_long_files(cat_E_files), dates_on_after)
    cat_K_shorts = sort_files_by_date(filter_short_files(cat_K_files), dates_on_after)
    cat_E_shorts = sort_files_by_date(filter_short_files(cat_E_files), dates_on_after)
    cat_K_13secs = sort_files_by_date(filter_13sec_short_files(cat_K_files), dates_on_after)
    cat_E_13secs = sort_files_by_date(filter_13sec_short_files(cat_E_files), dates_on_after)

    # display(cat_K_longs)
    # display(cat_K_shorts)
    # display(cat_K_13secs)
    # display(cat_E_longs)
    # display(cat_E_shorts)
    # display(cat_E_13secs)

    trans_list_of_K_files(cat_K_longs, cat_E_longs, conf_file)
    trans_list_of_K_files(cat_K_shorts, cat_E_shorts, conf_file)
    trans_list_of_K_files(cat_K_13secs, cat_E_13secs, conf_file)

    ist=IST(conf_file)
    for f in cat_K_longs:
        upload_a_ppt_ist(ist, f, 'LF')
    for f in cat_E_longs:
        upload_a_ppt_ist(ist, f, 'LF')
    for f in cat_K_shorts:
        upload_a_ppt_ist(ist, f, 'ST')
    for f in cat_E_shorts:
        upload_a_ppt_ist(ist, f, 'ST')
    for f in cat_K_13secs:
        upload_a_ppt_ist(ist, f, '13')
    for f in cat_E_13secs:
        upload_a_ppt_ist(ist, f, '13')


if __name__ == '__main__': 
    DATES_ON_AFTER = '2024-01-01' # None for today
    find_ppt_tranlate_and_upload(PPT_WORK_DIR, DATES_ON_AFTER, CONF_FILE)