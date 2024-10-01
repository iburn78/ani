#%% 
from pptx import Presentation
from datetime import datetime
import json 
import os

meta = {
    'ppt_path': 'data/ppt/',
    'ppt_file': '삼성전자_2024_2Q_E_2024-09-30.pptx',
}

max_size = 4500
slide_break = 2 
line_break = 0.7

txt_file_base = os.path.join(meta['ppt_path'], meta['ppt_file']).replace('.pptx', '')
txt_file_meta = os.path.join(meta['ppt_path'], f'meta.json')
ppt = Presentation(os.path.join(meta['ppt_path'], meta['ppt_file']))

header = '''<speak>
'''
footer = '''<break time="1s"/>
</speak>
'''

def write_to_file(content, current_file_number, current_size):
    txt_file = f"{txt_file_base}_{current_file_number}.txt"
    
    if current_size == 0:
        with open(txt_file, 'w', encoding='utf-8') as notes_file:
            notes_file.write(header)
            notes_file.write(content)
    else: 
        with open(txt_file, 'a', encoding='utf-8') as notes_file:
            notes_file.write(content)
    
    return current_size + len(content.encode('utf-8'))

file_number = 1
current_size = 0  
for slide_number, slide in enumerate(ppt.slides, start=0):
    if slide.notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text
        slide_content = f'.<mark name="slide{slide_number}"/>.\n<break time="{slide_break}s"/>\n'
        slide_content += notes.replace('\n', f'<break time="{line_break}s"/>\n') + f'<break time="{slide_break}s"/>\n'
    else:
        slide_content = f'.<mark name="slide{slide_number}"/>.\n<break time="{slide_break}s"/>\n'

    if current_size + len(slide_content.encode('utf-8')) > max_size:
        write_to_file(footer, file_number, current_size)
        file_number += 1
        current_size = 0

    current_size = write_to_file(slide_content, file_number, current_size)

write_to_file(footer, file_number, current_size)
with open(txt_file_meta, 'w', encoding='utf-8') as meta_file:
    meta['txt_file_number'] = file_number
    json.dump(meta, meta_file, ensure_ascii=False, indent=4)
