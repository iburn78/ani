#%% 
from pptx import Presentation
import pandas as pd
from PIL import Image
import os

class db_to_ppt:

    def __init__(self, )
pptx_filename = '삼성전자_K_2024-11-08_shorts_13sec.pptx'

working_dir = 'data/ppt/'
blank_filename = 'blank.pptx'
content_db_filename = 'content_db.xlsx'



def replace_shape_text(shape, new_text):
    if not shape.has_text_frame:
        raise Exception('Not a text shape')
    for paragraph in shape.text_frame.paragraphs:
        runs = paragraph.runs
        if runs:
            # Overwrite existing runs to preserve as much formatting as possible
            runs[0].text = new_text
            # Clear remaining runs if there are any extra
            for i in range(1, len(runs)):
                runs[i].text = ""
        else:
            # Fallback if no runs exist: add a new run with the new text
            new_run = paragraph.add_run()
            new_run.text = new_text
        if shape.text == new_text: 
            return 
        else: 
            raise Exception('Text replacement error')

def replace_shape_image(shape, new_image_path): 
    if hasattr(shape, "image"):
        # Get the position and size of the current image
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
                
        # Replace with a new image
        slide.shapes.add_picture(new_image_path, left, top, width=width, height=height)
        # Remove the old image shape
        slide.shapes._spTree.remove(shape._element)
    else: 
        raise Exception('Not an image shape')

def read_content_db(excel_path):
    # DB structure
    sheet_name = 'datasheet'
    preset_columns = ["name", "lang", "date", "suffix", "slide", "type", "title", "subtitle", "image_path", "image", "desc", "note"]
    
    if os.path.exists(excel_path):
        df = pd.read_excel(excel_path, sheet_name=sheet_name) 
    else:
        df = pd.DataFrame(columns=preset_columns)
        with pd.ExcelWriter(excel_path) as writer:
            df.to_excel(writer, sheet_name=sheet_name, index=False)
    return df

def list_slide_layouts(prs):
    for idx, layout in enumerate(prs.slide_layouts):
        print(f"Layout {idx}: {layout.name}")

def get_slide_type(prs, type):
    for idx, layout in enumerate(prs.slide_layouts):
        if type in layout.name: 
            return layout

def print_placeholder_idx(prs): 
    for slide_number, slide in enumerate(prs.slides, start=0):
        print(f"Slide number: {slide_number}")
        print('--------------')
        for placeholder in slide.placeholders:
            print(f"Placeholder Index: {placeholder.placeholder_format.idx}")
            print(f"Placeholder Type: {placeholder.placeholder_format.type}")
            if hasattr(placeholder, 'text'):
                print(f"Placeholder Text: '{placeholder.text}'")

def make_square_img(image_path, fill_color=(255, 255, 255, 0)):
    with Image.open(image_path) as img:
        width, height = img.size
        new_size = max(width, height)
        
        # Create a new image with the square size and fill color (transparent by default)
        new_img = Image.new("RGBA", (new_size, new_size), fill_color)
        # Paste the original image at the center of the new image
        new_img.paste(img, ((new_size - width) // 2, (new_size - height) // 2))
        new_img.save(image_path)


blank_path = os.path.join(working_dir, blank_filename)
content_db_path = os.path.join(working_dir, content_db_filename)
prs = Presentation(blank_path)
content_db = read_content_db(content_db_path)
content_db['date'] = content_db['date'].dt.strftime('%Y-%m-%d')
fa = pptx_filename.replace('.pptx', '').split('_')
target_db = content_db.loc[(content_db['name'] == fa[0])
            & (content_db['lang'] == fa[1])
            & (content_db['date'] == fa[2])
            & (content_db['suffix'] == fa[3]+'_'+fa[4])]

placeholder_idx_dict = {
    'title': {
        'date': 11, 
        'title': 0,
        'subtitle': None, 
        'image': None, 
        'desc': None, 
    }, 
    'image': {
        'date': None, 
        'title': 0, 
        'subtitle': 10, 
        'image': 11, 
        'desc': 12, 
    }, 
    'bullet': {
        'date': None, 
        'title': 0, 
        'subtitle': 10, 
        'image': None, 
        'desc': 12, 
    }, 
    'close': {
        'date': None, 
        'title': 0, 
        'subtitle': 10, 
        'image': None, 
        'desc': None, 
    }, 
}

for index, row in target_db.iterrows():
    slide = prs.slides.add_slide(get_slide_type(prs, row['type']))
    tdict = placeholder_idx_dict[row['type']]
    for i in ['date', 'title', 'subtitle', 'desc']: 
        if tdict[i] != None: 
            tph = next((obj for obj in slide.placeholders if obj.placeholder_format.idx == tdict[i]), None)
            tph.text = row[i]

    if tdict['image'] != None: 
        img = os.path.join(row['image_path'], row['image'])
        make_square_img(img)
        tph = next((obj for obj in slide.placeholders if obj.placeholder_format.idx == tdict['image']), None)
        tph.insert_picture(img)

        
prs.save(os.path.join(working_dir, pptx_filename))


