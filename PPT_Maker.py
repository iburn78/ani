#%% 
import os
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from PIL import Image
import math
from ani_tools import open_ppt_file, close_ppt_if_saved
import re

cd_ = os.path.dirname(os.path.abspath(__file__)) # .   
WORKING_DIR = os.path.join(cd_, 'data/ppt/')
OUTPUT_DIR = os.path.join(cd_, 'data/ppt/ppts/' )

# reading a content DB, and make a ppt file for a given v_id
# defining procedure and format within the class for a given blank template
class PPT_MAKER:
    BLANK_FILE_NAME = 'blank.pptx'
    CONTENT_DB_FILENAME = 'content_db.xlsx'
    CONTENT_DB_COLUMNS = ["v_id", "name", "lang", "date", "suffix", "slide", "type", "title", "subtitle", "image_path", "image", "desc", "note"] 
    CONTENT_DB_SHEETNAME = 'datasheet'
    SLIDE_TYPE = ['title', 'image', 'bullet', 'close']
    BULLET_PATTERN = r'^\[\s*(.*?)\s*\]$'  #  [whatever text] at the start of a line [, and end of a line ]
    SUPPLEMENT_PATTERN = r'\<\s*(.*?)\s*\>'  #  <whatever text>
    BULLET = '\u25A0'

    # CUSTOM SETTINGS
    PLACEHOLDER_IDX_DICT = {
        'title': {'date': 11, 'title': 0, 'subtitle': None, 'image': None, 'desc': None},
        'image': {'date': None, 'title': 0, 'subtitle': 10, 'image': 11, 'desc': 12},
        'bullet': {'date': None, 'title': 0, 'subtitle': 10, 'image': None, 'desc': 12, 'market': 13, 'rank': 14}, 
        'close': {'date': None, 'title': 0, 'subtitle': 10, 'image': None, 'desc': None},
    }
    FONT_SIZE_DICT_K = {
        'title': {'date': 40, 'title': 50, 'subtitle': None, 'image': None, 'desc': None},
        'image': {'date': None, 'title': 50, 'subtitle': 30, 'image': None, 'desc': 25},
        'bullet': {'date': None, 'title': 50, 'subtitle': 30, 'image': None, 'desc': 25},
        'close': {'date': None, 'title': 50, 'subtitle': 30, 'image': None, 'desc': None},
    }
    FONT_SIZE_DICT_E = {
        'title': {'date': 40, 'title': 42, 'subtitle': None, 'image': None, 'desc': None},
        'image': {'date': None, 'title': 42, 'subtitle': 27, 'image': None, 'desc': 22},
        'bullet': {'date': None, 'title': 42, 'subtitle': 27, 'image': None, 'desc': 22},
        'close': {'date': None, 'title': 42, 'subtitle': 27, 'image': None, 'desc': None},
    }
    EXT_FONT_REDUCTION = 3 # in bullet paragraph, font size above is for title in [ ]. detailed exp has EXT_FONT_REDUCTION amount less than this. Also usedin in ( ) too.

    COLOR_DICT = {
        "orange": (255, 165, 0),
        "dark_blue": (0, 0, 139),
        "green": (0, 128, 0),
        "dark_green": (0, 100, 0),
        "light_gray": (169, 169, 169),
        "dark_gray": (40, 40, 40),
    }

    INTRO_K = '13초컷 - 회사 하나 공부하자!'
    INTRO_E = '13sec cut - Learn a company!'
    INTRO_TEXTBOX_NAME = 'TextBox 2'

    def __init__(self, v_id=None, target_filename=None):
        self.prs = Presentation(PPT_MAKER.get_file_path(PPT_MAKER.BLANK_FILE_NAME))
        self.content_db = PPT_MAKER.read_content_db(PPT_MAKER.get_file_path(PPT_MAKER.CONTENT_DB_FILENAME))
        if v_id == None: 
            return 
        if self.get_target_db(v_id, target_filename):
            self.localizer()
            print(f'Target_db creation successful for {self.target_pptx_name} with length {len(self.target_db)}.')
        self.make_ppt()

    @staticmethod
    def get_file_path(filename):
        return os.path.join(WORKING_DIR, filename)

    @staticmethod
    def read_content_db(excel_path):
        # DB structure
        sheet_name = PPT_MAKER.CONTENT_DB_SHEETNAME
        preset_columns = PPT_MAKER.CONTENT_DB_COLUMNS

        if os.path.exists(excel_path):
            df = pd.read_excel(excel_path, sheet_name=sheet_name)
            df['date'] = pd.to_datetime(df['date'], errors='coerce').dt.strftime("%Y-%m-%d")
        else:
            df = pd.DataFrame(columns=preset_columns)
            with pd.ExcelWriter(excel_path) as writer:
                df.to_excel(writer, sheet_name=sheet_name, index=False)
        return df

    # you can either give v_id (unique) or a filename complying a strict format.
    def get_target_db(self, v_id, target_filename=None):
        if target_filename:
            target_db = self.get_target_db_by_filename(target_filename)
        else:
            target_db = self.content_db.loc[self.content_db['v_id'] == v_id]

        self.target_db = self.validate_target_db(v_id, target_filename, target_db)
        self.target_pptx_name = target_filename or self.generate_target_filename(target_db)
        self.validate_target_filename()
        return True

    def get_target_db_by_filename(self, target_filename):
        tf = target_filename.replace('.pptx', '').split('_')
        return self.content_db.loc[
            (self.content_db['name'] == tf[0]) &
            (self.content_db['lang'] == tf[1]) &
            (self.content_db['date'] == tf[2]) &
            (self.content_db['suffix'] == tf[3] + '_' + tf[4])
        ]

    def validate_target_db(self, v_id, target_filename, target_db):
        if len(target_db) == 0: 
            raise ValueError("No matching entry found")
        # Check if target_db is empty
        if not (target_db['v_id'] == v_id).all():
            raise ValueError(f"v_id ({v_id}) is not unique for the given filename: {target_filename}")
        # Check if 'name', 'lang', 'date', or 'suffix' columns have non-identical values
        if target_db['name'].nunique() > 1 or target_db['lang'].nunique() > 1 or target_db['date'].nunique() > 1 or target_db['suffix'].nunique() > 1:
            raise ValueError(f"Non-identical values found in 'name', 'lang', 'date', or 'suffix' columns for v_id: {v_id}")
        # Check if all values in 'slide' are integers
        if not target_db['slide'].apply(lambda x: isinstance(x, int)).all():
            raise ValueError(f"All values in 'slide' column must be integers for v_id: {v_id}")
        # Check if all values in 'slide' are different
        if target_db['slide'].nunique() != len(target_db):
            raise ValueError(f"Duplicate values found in 'slide' column for v_id: {v_id}")
        return target_db

    def generate_target_filename(self, target_db):
        return f"{target_db['name'].iloc[0]}_{target_db['lang'].iloc[0]}_{target_db['date'].iloc[0]}_{target_db['suffix'].iloc[0]}.pptx"

    def validate_target_filename(self):
        if self.target_pptx_name == PPT_MAKER.BLANK_FILE_NAME:
            raise Exception('Check target file name: should not be the blank file name')

    def make_ppt(self):
        for index, row in self.target_db.iterrows():
            slide = self.prs.slides.add_slide(self._get_slide_type(row['type']))
            self.populate_slide_with_data(slide, row)
            self.add_image_to_slide(slide, row)
            self.set_note(slide, row)

        close_ppt_if_saved(self.target_pptx_name)
        self.final_ppt_path_filename = os.path.join(OUTPUT_DIR, self.target_pptx_name)
        self.prs.save(self.final_ppt_path_filename)
        print("PPT save completed ---- ")

    def populate_slide_with_data(self, slide, row):
        tdict = PPT_MAKER.PLACEHOLDER_IDX_DICT[row['type']]
        for i in ['date', 'title', 'subtitle', 'desc']:
            if tdict[i] is None: continue # should check None explictly, cause value 0 is used too!
            tph = next((obj for obj in slide.placeholders if obj.placeholder_format.idx == tdict[i]), None)
            if tph and tph.text_frame:
                value = row[i]
                # Check for NaN and None explicitly
                if value is None or (isinstance(value, str) and value.replace('\n','').strip() == "") or (isinstance(value, float) and math.isnan(value)): 
                    slide.shapes._spTree.remove(tph._element)
                    continue
                value = PPT_MAKER.text_preprocessor(value)

                text_frame = tph.text_frame
                text_frame.clear()  # Clear the existing text
                # text wrapping works only for English words
                text_frame.word_wrap = True
                # first paragraph is always given
                paragraph = text_frame.paragraphs[0]  # Get the first paragraph

                if self.target_db['lang'].iloc[0] == 'K':
                    font_size_pt = PPT_MAKER.FONT_SIZE_DICT_K[row['type']][i]
                else: # self.target_db['lang'].iloc[0] == 'K':
                    font_size_pt = PPT_MAKER.FONT_SIZE_DICT_E[row['type']][i]

                # BULLET slide, DESC paragraph has formatting option with [ ]
                if row['type'] == 'bullet' and i == 'desc':
                    for line in value.splitlines():
                        run = paragraph.add_run() 
                        match = re.match(PPT_MAKER.BULLET_PATTERN, line) # match only from the start of the line
                        if match:
                            run.text=self.text_wrapper(PPT_MAKER.BULLET + ' ' +match.group(1), font_size_pt, tph.width)+'\n'
                            run.font.size = PPT_MAKER.pt_to_emu(font_size_pt)
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(*PPT_MAKER.COLOR_DICT['dark_blue'])
                        else: 
                            run.text=self.text_wrapper(line, font_size_pt-PPT_MAKER.EXT_FONT_REDUCTION, tph.width)+'\n'
                            run.font.size = PPT_MAKER.pt_to_emu(font_size_pt-PPT_MAKER.EXT_FONT_REDUCTION)
                            run.font.bold = False
                else: 
                    # handling supplement pattern < > and give it different style
                    search = re.search(PPT_MAKER.SUPPLEMENT_PATTERN, value)
                    value = re.sub(PPT_MAKER.SUPPLEMENT_PATTERN, "", value)
                    run = paragraph.add_run()  # Get the first run of the paragraph
                    text_wrapped = self.text_wrapper(value, font_size_pt, tph.width) 
                    run.text = text_wrapped
                    run.font.size = PPT_MAKER.pt_to_emu(font_size_pt)
                    if search:
                        run = paragraph.add_run()  
                        if self.target_db['lang'].iloc[0] != 'K' or len(text_wrapped.splitlines()) != 1:
                            run.text = ' ('+search.group(1)+')'
                        else: 
                            run.text =  '\n('+search.group(1)+')'
                        run.font.size = PPT_MAKER.pt_to_emu(font_size_pt-PPT_MAKER.EXT_FONT_REDUCTION)
                        run.font.color.rgb = RGBColor(*PPT_MAKER.COLOR_DICT['dark_gray'])
                        run.font.bold = False

    @staticmethod
    def pt_to_emu(pt):
        return int(pt * 12700) # 914400 EMUs per inch, 72 points per inch

    @staticmethod
    def text_preprocessor(text):
        # Preprocessing 
        # Add newlines around [ ] if needed and apply strip for each line
        text = str(text).strip()
        text = re.sub(rf'({PPT_MAKER.BULLET_PATTERN})(?!\n)', r'\1\n', text, flags=re.MULTILINE)
        text = re.sub(rf'(?<!\n)({PPT_MAKER.BULLET_PATTERN})', r'\n\1', text, flags=re.MULTILINE)
        # Apply strip() each line
        text = '\n'.join(line.strip() for line in text.splitlines()).strip()
        # Replace 3 or more newlines with 2
        text = re.sub(r'\n{3,}', '\n\n', text)   
        # If any line is in "" or '', remove them
        text = '\n'.join(
        line[1:-1] if (line.startswith('"') and line.endswith('"')) or (line.startswith("'") and line.endswith("'"))
        else line
        for line in text.splitlines()
        )

        return text

    # text_wrapper wraps words in case if Korean 
    # In English case, it is automatically done by ppt. However, below applied in English as well for consistency
    def text_wrapper(self, text, fontsize, width): # fontsize is in pt (height), width is in EMU
        if self.target_db['lang'].iloc[0] == 'K':
            FONT_WIDTH_ADJUSTMENT = 0.82  # Apply []% to the height for width calculation
            TEXTBOX_WIDTH_MARGIN = 0.85
        else: 
            # return text # if you don't want to apply text wrapper in English case
            FONT_WIDTH_ADJUSTMENT = 0.77  # Apply []% to the height for width calculation
            TEXTBOX_WIDTH_MARGIN = 0.85

        bullet_to_font_size = 1.2 
        Eng_to_Kor_size = 0.67
        space_to_font_size = 0.25
        # Font size is Height of the font, not width
        fontsize = FONT_WIDTH_ADJUSTMENT*fontsize # Adjust the font size to fit the width
        width_in_pt = (width / 914400) * 72
        max_line_width = width_in_pt * TEXTBOX_WIDTH_MARGIN  # Apply []% of the width for text

        lines = []
    
        for pg in text.splitlines():  # pg stands for paragraph
            words = pg.split()
            current_line = ""
            current_length = 0  # Track the length of the current line in points
        
            for word in words:
                # Calculate the word length in points
                word_length = sum(fontsize if '\uAC00' <= char <= '\uD7A3' else fontsize*bullet_to_font_size if char ==  PPT_MAKER.BULLET else fontsize*Eng_to_Kor_size for char in word) + fontsize*space_to_font_size  # Add space width
            
                # Check if adding the word exceeds the max width
                if current_length + word_length > max_line_width:
                    # Append the current line and reset
                    lines.append(current_line.strip())
                    current_line = word
                    current_length = word_length
                else:
                    # Add the word to the current line
                    current_line += " " + word
                    current_length += word_length
        
            # Add the last line of the paragraph
            lines.append(current_line.strip())

        return "\n".join(lines).strip()

    def add_image_to_slide(self, slide, row):
        tdict = PPT_MAKER.PLACEHOLDER_IDX_DICT[row['type']]
        if tdict['image'] is not None: # should not None explictly!
            img_path = os.path.join(row['image_path'], row['image'])
            self.make_img_square(img_path)
            tph = next((obj for obj in slide.placeholders if obj.placeholder_format.idx == tdict['image']), None)
            if tph:
                tph.insert_picture(img_path)

    def make_img_square(self, image_path, fill_color=(255, 255, 255, 0)):
        with Image.open(image_path) as img:
            width, height = img.size
            new_size = max(width, height)
            new_img = Image.new("RGBA", (new_size, new_size), fill_color)
            new_img.paste(img, ((new_size - width) // 2, (new_size - height) // 2))
            new_img.save(image_path)

    def set_note(self, slide, row): 
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            value = row['note']
            if value is None or (isinstance(value, str) and value.replace('\n','').strip() == "") or (isinstance(value, float) and math.isnan(value)): 
                value = ''
            slide.notes_slide.notes_text_frame.text = str(value).strip()

    def localizer(self):
        title_layout = self._get_slide_type('title')
        # Change the intro 
        hd = [s for s in title_layout.shapes if s.name == PPT_MAKER.INTRO_TEXTBOX_NAME][0]
    
        # Get the first paragraph and the first run
        # this is to keep the formating of the textbox
        tg_run = hd.text_frame.paragraphs[0].runs[0]

        if self.target_db['lang'].iloc[0] == 'K':
            tg_run.text = PPT_MAKER.INTRO_K
        else: 
            tg_run.text = PPT_MAKER.INTRO_E
        ### ADD MORE LOCALIZE SET UPS HERE ###

    def _get_slide_type(self, type):
        for idx, layout in enumerate(self.prs.slide_layouts):
            if type in layout.name: 
                return layout

    def list_slide_layouts(self):
        for idx, layout in enumerate(self.prs.slide_layouts):
            print(f"Layout {idx}: {layout.name}")

    def print_placeholder_idx(self): 
        for slide_number, slide in enumerate(self.prs.slides, start=0):
            print(f"Slide number: {slide_number}")
            print('--------------')
            for placeholder in slide.placeholders:
                print(f"Placeholder Index: {placeholder.placeholder_format.idx}")
                print(f"Placeholder Type: {placeholder.placeholder_format.type}")
                if hasattr(placeholder, 'text'):
                    print(f"Placeholder Text: '{placeholder.text}'")

    @staticmethod
    def replace_shape_text(shape, new_text):
        if not shape.has_text_frame:
            raise Exception('Not a text shape')
        for paragraph in shape.text_frame.paragraphs:
            runs = paragraph.runs
            if runs:
                # Overwrite existing runs to preserve as much formatting as possible
                runs[0].text = new_text
                # Clear remaining runs if there are any extra
                for i in range(1, len(runs)):
                    runs[i].text = ""
            else:
                # Fallback if no runs exist: add a new run with the new text
                new_run = paragraph.add_run()
                new_run.text = new_text
            if shape.text == new_text: 
                return True
            else: 
                raise Exception('Text replacement error')

    @staticmethod
    def replace_shape_image(shape, new_image_path): 
        if hasattr(shape, "image"):
            # Get the position and size of the current image
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
                
            # Remove the old image shape
            shape._spTree.remove(shape._element)
            # Replace with a new image
            shape.add_picture(new_image_path, left, top, width=width, height=height)
        else: 
            raise Exception('Not an image shape')

if __name__ == '__main__': 
    v_id = 2
    pm = PPT_MAKER(v_id)
    open_ppt_file(pm.final_ppt_path_filename)