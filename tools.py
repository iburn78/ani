#%%
from pptx import Presentation
from google.cloud import texttospeech_v1beta1 as tts
from moviepy.editor import ImageClip, concatenate_videoclips, AudioFileClip
import os
# import json 

def ppt_to_video(meta): 
    max_size = 4500
    slide_break = 2 
    line_break = 0.7
    fps = 24

    meta['ppt_path'] = 'data/ppt/'
    meta['voice_path'] = 'data/voice/'
    meta['image_prefix'] = '슬라이드'

    meta_updated = ppt_to_text(meta, max_size=max_size, slide_break=slide_break, line_break=line_break)
    timepoints = ppt_tts(meta_updated)
    create_video_from_ppt_and_voice(meta_updated, timepoints=timepoints, fps=fps)

def ppt_to_text(meta, max_size = 4500, slide_break = 2, line_break = 0.7):
    header = '''<speak>
    '''
    footer = '''<break time="1s"/>
    </speak>
    '''
    
    file_number = 0
    current_size = _write_to_file(header, file_number, 0, meta)
    ppt = Presentation(os.path.join(meta['ppt_path'], meta['ppt_file']))
    
    for slide_number, slide in enumerate(ppt.slides):
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            slide_content = f'.<mark name="slide{slide_number}"/>.\n<break time="{round(slide_break/2,1)}s"/>\n'
            slide_content += notes.replace('\n', f'<break time="{line_break}s"/>\n') + f'<break time="{slide_break}s"/>\n'
        else:
            slide_content = f'.<mark name="slide{slide_number}"/>.\n<break time="{slide_break}s"/>\n'

        if current_size + len(slide_content.encode('utf-8')) > max_size:
            _write_to_file(footer, file_number, current_size, meta)
            file_number += 1
            current_size = 0
            slide_content = header + slide_content

        current_size = _write_to_file(slide_content, file_number, current_size, meta)

    _write_to_file(footer, file_number, current_size, meta)

    meta_file = os.path.join(meta['ppt_path'], 'meta.json')
    meta['txt_file_number'] = file_number+1
    # with open(meta_file, 'w', encoding='utf-8') as meta_json:
    #     json.dump(meta, meta_json, ensure_ascii=False, indent=4)

    return meta

def _write_to_file(content, current_file_number, current_size, meta):
    txt_file = f"{os.path.join(meta['voice_path'], meta['ppt_file'].replace('.pptx', ''))}_{current_file_number}.txt"

    mode = 'w' if current_size == 0 else 'a'
    with open(txt_file, mode, encoding='utf-8') as notes_file:
        notes_file.write(content)
    
    return current_size + len(content.encode('utf-8'))

def ppt_tts(meta):
    ppt_file = os.path.join(meta['ppt_path'], meta['ppt_file'])
    os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = '../config/google_cloud.json'

    client = tts.TextToSpeechClient()
    language_code = 'ko-KR' if meta['lang'] == 'K' else 'en-US' 
    tag = 'D' if meta['lang'] == 'K' else 'B' 
    name = language_code+'-Wavenet-'+tag 
    if meta['wave'] == True: # WaveNet voice (1 mil words/month vs 4 mil in basic)
        voice = tts.VoiceSelectionParams(language_code=language_code, name=name, ssml_gender=tts.SsmlVoiceGender.MALE)
    else:
        voice = tts.VoiceSelectionParams(language_code=language_code, ssml_gender=tts.SsmlVoiceGender.MALE)
    audio_config = tts.AudioConfig(audio_encoding=tts.AudioEncoding.MP3)
    
    timepoint_dict = {}
    for i in range(meta['txt_file_number']):
        txt_file = f"{os.path.join(meta['voice_path'], meta['ppt_file'].replace('.pptx', '_'+str(i)+'.txt'))}"
        voice_file =ppt_file.replace('/ppt/','/voice/').replace('.pptx', '_'+str(i)+'.mp3')

        with open(txt_file, 'r', encoding='utf-8') as file:
            text_content = file.read()

        synthesis_input = tts.SynthesisInput(ssml=text_content)
        request = tts.SynthesizeSpeechRequest(
            input=synthesis_input,
            voice=voice,
            audio_config=audio_config, 
            enable_time_pointing=[tts.SynthesizeSpeechRequest.TimepointType.SSML_MARK]
        )
        response = client.synthesize_speech(request=request)

        with open(voice_file, "wb") as out:
            out.write(response.audio_content)
            print(voice_file + ' done')

        timepoint_list = []
        if response.timepoints:
            for time_point in response.timepoints:
                print(f'Mark name: {time_point.mark_name}, Time: {time_point.time_seconds} seconds')
                timepoint_list.append([int(time_point.mark_name[5:]), time_point.time_seconds])
        else:
            print('No timepoints found.')
        timepoint_dict[voice_file] = timepoint_list

    # tps_file = os.path.join(meta['voice_path'], 'timepoints.json')
    # with open(tps_file, 'w', encoding='utf-8') as timepoints_json:
    #     json.dump(timepoint_dict, timepoints_json, ensure_ascii=False, indent=4)

    return timepoint_dict

def create_video_from_ppt_and_voice(meta, timepoints, fps=24):
    images_folder = os.path.join(meta['ppt_path'], meta['ppt_file'].replace('.pptx',''))
    output_file = os.path.join(meta['ppt_path'], meta['ppt_file'].replace('.pptx', '.mp4'))
    image_prefix = meta['image_prefix']
    video_with_audios = []

    for audio_file, slide_times in timepoints.items():
        audio_clip = AudioFileClip(audio_file)

        video_clips = []
        for i in range(len(slide_times)):
            start_time = slide_times[i][1]  # Get the start time for the slide
            if i < len(slide_times)-1:
                end_time = slide_times[i + 1][1]  # Get the end time for the next slide
            else:
                end_time = audio_clip.duration
            slide_number = slide_times[i][0]

            # Construct the image filename
            slide_image_filename = f'{image_prefix}{slide_number}.PNG'
            slide_image_path = os.path.join(images_folder, slide_image_filename)

            # Load the slide image
            slide_clip = ImageClip(slide_image_path).set_duration(end_time - start_time).set_start(start_time)
            video_clips.append(slide_clip)

        # Concatenate video clips for the current audio
        video_with_audio = concatenate_videoclips(video_clips)
        video_with_audio = video_with_audio.set_audio(audio_clip).volumex(2)
        video_with_audios.append(video_with_audio)

    # Concatenate all videos into one final video
    final_video = concatenate_videoclips(video_with_audios)

    # Set fps for the final video
    final_video.fps = fps
    
    # final_video.write_videofile(output_file, codec="libx264")
    final_video.write_videofile(
        output_file,
        codec="libx264",
    )
    print('video generation done')